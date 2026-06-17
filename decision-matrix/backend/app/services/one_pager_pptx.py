"""Generate PowerPoint one-pager export (FR-11.2.2)."""

from __future__ import annotations

import base64
import io
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from app.geo.constants import SUBTYPE_LABELS as _MANIFEST_SUBTYPE_LABELS
from app.models import OnePager
from app.services.one_pager_builder import ENG_LABELS

# Brand palette (aligned with one-pager web UI)
C_PRIMARY = RGBColor(0x0B, 0x5C, 0xAD)
C_PRIMARY_LIGHT = RGBColor(0xE8, 0xF1, 0xFA)
C_ACCENT = RGBColor(0x00, 0xA8, 0x96)
C_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
C_MUTED = RGBColor(0x5A, 0x6D, 0x85)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BORDER = RGBColor(0xE0, 0xE4, 0xEA)
C_OK_BG = RGBColor(0xE8, 0xF5, 0xE9)
C_WARN_BG = RGBColor(0xFF, 0xEB, 0xEE)
C_PANEL = RGBColor(0xF4, 0xF6, 0xF9)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

SUBTYPE_LABELS: dict[str, str] = {
    **_MANIFEST_SUBTYPE_LABELS,
    "gtes": "ГТЭС/ГПЭС",
    "pads": "Кустовые площадки",
}

STATUS_LABELS: dict[str, str] = {
    "within_limit": "В пределах",
    "exceeds_limit": "Превышение",
    "not_required": "Не требуется",
    "construction_required": "Строительство",
    "computed": "Расчёт",
}

ENG_KEYS: list[tuple[str, str]] = [
    ("eng_power", "Электроснабжение"),
    ("eng_injection", "Закачка"),
    ("eng_gas", "Газ"),
    ("eng_oil_preparation", "Подготовка нефти"),
    ("eng_well_gathering", "Сбор скважин"),
    ("eng_transport", "Транспорт"),
]

GANTT_COLORS = [
    RGBColor(0x0B, 0x5C, 0xAD),
    RGBColor(0x12, 0x76, 0xC7),
    RGBColor(0x1A, 0x8F, 0xE0),
    RGBColor(0x23, 0xA8, 0xF5),
    RGBColor(0x2F, 0xC0, 0xFF),
    RGBColor(0x00, 0xA8, 0x96),
]


def _exports_dir() -> Path:
    from app.core.paths import data_dir

    return data_dir("exports")


def _fmt_font(run, *, size: int = 10, bold: bool = False, color: RGBColor | None = None) -> None:
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def _set_para(
    text_frame,
    text: str,
    *,
    size: int = 10,
    bold: bool = False,
    color: RGBColor | None = None,
    align: PP_ALIGN | None = None,
) -> None:
    text_frame.clear()
    p = text_frame.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    _fmt_font(run, size=size, bold=bold, color=color)


def _add_rect(slide, left, top, width, height, fill: RGBColor, *, line: RGBColor | None = None, line_w: float = 0.75):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    return shape


def _roadmap_to_gantt(stages: list[dict[str, Any]]) -> list[dict[str, Any]]:
    finite = [s.get("duration_months") for s in stages if isinstance(s.get("duration_months"), (int, float))]
    cumulative = len(finite) >= 2 and all(finite[i] >= finite[i - 1] for i in range(1, len(finite)))
    cursor = 0
    out: list[dict[str, Any]] = []
    for i, step in enumerate(stages):
        dur = step.get("duration_months")
        if dur is None:
            out.append(
                {"stage": step.get("stage", ""), "start": cursor, "end": None, "open": True, "index": i}
            )
            continue
        end = int(dur)
        if cumulative:
            start = 0 if i == 0 else int(stages[i - 1].get("duration_months") or 0)
            cursor = end
        else:
            start = cursor
            cursor += end
        out.append(
            {
                "stage": step.get("stage", ""),
                "start": start,
                "end": end,
                "open": False,
                "index": i,
                "duration": end - start if cumulative else end,
            }
        )
    return out


def _gantt_span(segments: list[dict[str, Any]]) -> int:
    ends = [s["end"] for s in segments if s.get("end") is not None]
    return max(max(ends) if ends else 36, 12)


def _group_rows(rows: list[dict[str, Any]]) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    internal: list[dict[str, Any]] = []
    external: list[dict[str, Any]] = []
    line_subtypes = {"autoroad", "oil_pipeline", "water_pipeline", "power_line", "gas_pipeline"}
    for row in rows:
        pt = str(row.get("param_type", ""))
        sub = str(row.get("subtype", ""))
        if pt in ("external", "external_linear") or sub in line_subtypes:
            external.append(row)
        else:
            internal.append(row)
    return internal, external


def _row_line(row: dict[str, Any]) -> tuple[str, str, str, str]:
    subtype = SUBTYPE_LABELS.get(str(row.get("subtype", "")), str(row.get("subtype", "")))
    name = row.get("object_name")
    label = f"{subtype}" + (f" · {name}" if name else "")
    dist = row.get("distance_km")
    cost = row.get("cost_mln")
    dist_s = f"{dist:.1f} км" if isinstance(dist, (int, float)) else "—"
    cost_s = f"{cost:.1f} млн ₽" if isinstance(cost, (int, float)) else "—"
    status = STATUS_LABELS.get(str(row.get("status", "")), str(row.get("status", "")))
    return label, dist_s, cost_s, status


def _status_fill(status: str) -> RGBColor | None:
    if status == "exceeds_limit":
        return C_WARN_BG
    if status in ("within_limit", "computed"):
        return C_OK_BG
    return None


def _add_header(slide, one_pager: OnePager, poi_name: str) -> None:
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.52), C_PRIMARY)
    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.08), Inches(9.5), Inches(0.38))
    _set_para(title_box.text_frame, one_pager.title, size=20, bold=True, color=C_WHITE)

    date_str = str(one_pager.report_date or "")
    date_box = slide.shapes.add_textbox(Inches(10.2), Inches(0.1), Inches(2.8), Inches(0.34))
    tf = date_box.text_frame
    _set_para(tf, date_str, size=11, color=C_WHITE, align=PP_ALIGN.RIGHT)

    meta_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.56), Inches(12.4), Inches(0.22))
    parts = [
        p
        for p in [
            f"Координаты: {one_pager.coordinates}" if one_pager.coordinates else None,
            one_pager.engineer_name,
            f"Участок: {poi_name}" if poi_name else None,
        ]
        if p
    ]
    _set_para(meta_box.text_frame, " · ".join(parts), size=9, color=C_MUTED)


def _add_map_panel(slide, map_b64: str | None) -> None:
    left, top, w, h = Inches(0.35), Inches(0.88), Inches(4.85), Inches(3.05)
    _add_rect(slide, left, top, w, h, C_PANEL, line=C_BORDER)
    inner_l = left + Inches(0.06)
    inner_t = top + Inches(0.06)
    inner_w = w - Inches(0.12)
    inner_h = h - Inches(0.12)

    if map_b64:
        try:
            raw = map_b64.split(",", 1)[-1] if "," in map_b64 else map_b64
            data = base64.b64decode(raw)
            slide.shapes.add_picture(io.BytesIO(data), inner_l, inner_t, width=inner_w, height=inner_h)
            return
        except Exception:
            pass

    ph = slide.shapes.add_textbox(inner_l, inner_t + Inches(1.2), inner_w, Inches(0.5))
    _set_para(ph.text_frame, "Карта недоступна", size=12, color=C_MUTED, align=PP_ALIGN.CENTER)


def _fill_table_cell(cell, text: str, *, bold: bool = False, size: int = 8, fill: RGBColor | None = None, color: RGBColor | None = None):
    cell.text = text
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0] if p.runs else p.add_run()
    run.text = text
    _fmt_font(run, size=size, bold=bold, color=color or C_TEXT)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill


def _add_variant_table(slide, final_data: dict[str, Any], poi_name: str) -> None:
    left, top, width = Inches(5.35), Inches(0.88), Inches(7.55)
    title_box = slide.shapes.add_textbox(left, top, width, Inches(0.28))
    _set_para(
        title_box.text_frame,
        f"Анализ окружения: {poi_name}" if poi_name else "Анализ окружения",
        size=12,
        bold=True,
        color=C_PRIMARY,
    )

    rows_data = final_data.get("analysis_rows") or []
    internal, external = _group_rows(rows_data)

    table_rows: list[tuple[str, str, str, str, str, str | None]] = []
    for section, items in (("Внутренние", internal), ("Внешние", external)):
        if not items:
            table_rows.append((section, "—", "—", "—", "—", None))
            continue
        for row in items:
            label, dist, cost, status = _row_line(row)
            table_rows.append((section, label, dist, cost, status, str(row.get("status", ""))))

    n = len(table_rows) + 1
    table_shape = slide.shapes.add_table(n, 5, left, top + Inches(0.32), width, Inches(min(2.55, 0.22 * n)))
    table = table_shape.table
    col_w = [Inches(1.05), Inches(2.55), Inches(0.85), Inches(1.05), Inches(1.05)]
    for i, w in enumerate(col_w):
        table.columns[i].width = w

    headers = ("Группа", "Подтип", "Дист.", "Стоимость", "Статус")
    for ci, h in enumerate(headers):
        _fill_table_cell(table.cell(0, ci), h, bold=True, size=8, fill=C_PRIMARY_LIGHT, color=C_PRIMARY)

    prev_section = ""
    for ri, (section, label, dist, cost, status, raw_status) in enumerate(table_rows, start=1):
        sec = "" if section == prev_section else section
        prev_section = section
        row_fill = _status_fill(raw_status or "")
        _fill_table_cell(table.cell(ri, 0), sec, size=7, color=C_MUTED)
        _fill_table_cell(table.cell(ri, 1), label, size=8, fill=row_fill)
        _fill_table_cell(table.cell(ri, 2), dist, size=8, fill=row_fill)
        _fill_table_cell(table.cell(ri, 3), cost, size=8, fill=row_fill)
        _fill_table_cell(table.cell(ri, 4), status, size=7, fill=row_fill)


def _add_total_bar(slide, total: Any) -> None:
    left, top, width, height = Inches(5.35), Inches(3.52), Inches(7.55), Inches(0.36)
    _add_rect(slide, left, top, width, height, C_PRIMARY_LIGHT, line=C_PRIMARY, line_w=1)
    box = slide.shapes.add_textbox(left + Inches(0.12), top + Inches(0.05), width - Inches(0.2), height)
    total_s = f"{total:.1f} млн ₽" if isinstance(total, (int, float)) else "—"
    _set_para(box.text_frame, f"ИТОГО: {total_s}", size=14, bold=True, color=C_PRIMARY)


def _eng_label(key: str, value: str) -> str:
    return ENG_LABELS.get(key, {}).get(value, value)


def _add_eng_section(slide, eng: dict[str, Any], equipment_mln: Any) -> None:
    top = Inches(3.98)
    title = slide.shapes.add_textbox(Inches(0.35), top, Inches(12.6), Inches(0.22))
    _set_para(title.text_frame, "Инженерные параметры", size=11, bold=True, color=C_TEXT)

    col_count = 3
    pill_w = Inches(3.95)
    pill_h = Inches(0.26)
    gap_x = Inches(0.1)
    gap_y = Inches(0.06)
    base_x = Inches(0.35)
    base_y = top + Inches(0.24)

    for idx, (key, _title) in enumerate(ENG_KEYS):
        col = idx % col_count
        row = idx // col_count
        x = base_x + col * (pill_w + gap_x)
        y = base_y + row * (pill_h + gap_y)
        val = str(eng.get(key, "") or "—")
        label = _eng_label(key, val) if val != "—" else "—"
        pill = _add_rect(slide, x, y, pill_w, pill_h, C_PRIMARY_LIGHT, line=C_PRIMARY, line_w=0.5)
        tf = pill.text_frame
        tf.margin_left = Pt(6)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _set_para(tf, label, size=7, color=C_PRIMARY)

    if isinstance(equipment_mln, (int, float)):
        eq_y = base_y + 2 * (pill_h + gap_y) + Inches(0.04)
        eq = slide.shapes.add_textbox(Inches(0.35), eq_y, Inches(8), Inches(0.2))
        _set_para(
            eq.text_frame,
            f"Стоимость инженерного оборудования: {equipment_mln:.1f} млн ₽",
            size=9,
            color=C_MUTED,
        )


def _add_gantt(slide, roadmap: list[dict[str, Any]]) -> None:
    if not roadmap:
        return
    segments = _roadmap_to_gantt(roadmap)
    span = _gantt_span(segments)

    top = Inches(4.72)
    title = slide.shapes.add_textbox(Inches(0.35), top, Inches(12.6), Inches(0.22))
    _set_para(title.text_frame, "Дорожная карта", size=11, bold=True, color=C_TEXT)

    label_w = Inches(1.15)
    chart_l = Inches(0.35) + label_w + Inches(0.08)
    chart_w = Inches(11.0)
    row_h = Inches(0.22)
    track_top = top + Inches(0.28)

    for si, seg in enumerate(segments):
        row_top = track_top + si * (row_h + Inches(0.04))
        lbl = slide.shapes.add_textbox(Inches(0.35), row_top, label_w, row_h)
        _set_para(lbl.text_frame, str(seg.get("stage", "")), size=7, bold=True, color=C_TEXT)

        _add_rect(slide, chart_l, row_top, chart_w, row_h, C_PANEL, line=C_BORDER, line_w=0.5)

        start = seg.get("start", 0)
        if seg.get("open"):
            bar_l = chart_l + chart_w * (start / span)
            bar_w = chart_w - (bar_l - chart_l)
            bar = _add_rect(slide, bar_l, row_top + Inches(0.03), bar_w, row_h - Inches(0.06), C_ACCENT)
        elif seg.get("end") is not None:
            end = seg["end"]
            bar_l = chart_l + chart_w * (start / span)
            bar_w = chart_w * ((end - start) / span)
            color = GANTT_COLORS[seg.get("index", 0) % len(GANTT_COLORS)]
            bar = _add_rect(slide, bar_l, row_top + Inches(0.03), max(bar_w, Inches(0.08)), row_h - Inches(0.06), color)
            dur = seg.get("duration")
            if isinstance(dur, (int, float)) and dur > 0 and bar_w > Inches(0.35):
                tf = bar.text_frame
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                _set_para(tf, f"{int(dur)} мес.", size=6, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    axis_top = track_top + len(segments) * (row_h + Inches(0.04)) + Inches(0.02)
    axis = slide.shapes.add_textbox(chart_l, axis_top, chart_w, Inches(0.18))
    ticks = [0]
    step = 6 if span <= 24 else 12
    ticks.extend(list(range(step, span, step)))
    if span not in ticks:
        ticks.append(span)
    _set_para(axis.text_frame, " — ".join(str(t) for t in ticks) + " мес.", size=6, color=C_MUTED)


def _add_recommendation(slide, text: str) -> None:
    if not text:
        return
    left, top, width, height = Inches(0.35), Inches(6.05), Inches(12.55), Inches(1.15)
    _add_rect(slide, left, top, width, height, C_PANEL, line=C_BORDER)
    accent = _add_rect(slide, left, top, Inches(0.06), height, C_PRIMARY)
    accent.line.fill.background()

    box = slide.shapes.add_textbox(left + Inches(0.14), top + Inches(0.08), width - Inches(0.22), height - Inches(0.12))
    tf = box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = "Рекомендация"
    _fmt_font(r1, size=10, bold=True, color=C_PRIMARY)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = text
    _fmt_font(r2, size=9, color=C_TEXT)


def _add_watermark(slide) -> None:
    box = slide.shapes.add_textbox(Inches(10.5), Inches(6.85), Inches(2.5), Inches(0.35))
    _set_para(box.text_frame, "СППР", size=18, bold=True, color=RGBColor(0xE8, 0xEC, 0xF0))


def generate_one_pager_pptx(one_pager: OnePager) -> Path:
    """Build PPTX file on disk; return absolute path."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    final_data = one_pager.final_variant_data or {}
    eng = one_pager.engineering_params or {}
    total = final_data.get("total_cost_mln")
    equipment = final_data.get("equipment_cost_mln")
    poi_name = str(final_data.get("poi_name", "") or one_pager.title.split(" — ")[-1])

    _add_header(slide, one_pager, poi_name)
    _add_map_panel(slide, one_pager.map_snapshot_base64)
    _add_variant_table(slide, final_data, poi_name)
    _add_total_bar(slide, total)
    _add_eng_section(slide, eng, equipment)
    _add_gantt(slide, one_pager.roadmap or [])
    _add_recommendation(slide, one_pager.recommendation_text or "")
    _add_watermark(slide)

    out_dir = _exports_dir() / str(one_pager.project_id)
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / f"{one_pager.id}.pptx"
    prs.save(str(out_path))
    return out_path
